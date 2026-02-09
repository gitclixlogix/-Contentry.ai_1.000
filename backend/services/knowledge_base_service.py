"""
Knowledge Base Service for Strategic Profiles
Handles document processing, chunking, embedding, and RAG queries using ChromaDB.
"""

import os
import logging
import asyncio
from typing import List, Dict, Optional, Any
from uuid import uuid4
from datetime import datetime, timezone
from pathlib import Path
import hashlib

# Document processing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# Vector database
import chromadb
from chromadb.config import Settings

logger = logging.getLogger(__name__)

# Configuration
CHUNK_SIZE = 500  # characters per chunk
CHUNK_OVERLAP = 50  # overlap between chunks
UPLOADS_DIR = Path(__file__).parent.parent / "uploads" / "knowledge_base"
CHROMADB_DIR = Path(__file__).parent.parent / "data" / "chromadb"

# Ensure directories exist
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
CHROMADB_DIR.mkdir(parents=True, exist_ok=True)


class KnowledgeBaseService:
    """Service for managing knowledge base documents and RAG queries."""
    
    _instance = None
    _chroma_client = None
    _embedding_function = None
    
    def __init__(self, db=None):
        self.db = db
        self._init_chromadb()
    
    def _init_chromadb(self):
        """Initialize ChromaDB client."""
        if KnowledgeBaseService._chroma_client is None:
            logger.info(f"Initializing ChromaDB at {CHROMADB_DIR}")
            KnowledgeBaseService._chroma_client = chromadb.PersistentClient(
                path=str(CHROMADB_DIR),
                settings=Settings(
                    anonymized_telemetry=False,
                    allow_reset=True
                )
            )
            logger.info("ChromaDB initialized successfully")
        
        self.client = KnowledgeBaseService._chroma_client
    
    def _get_collection(self, profile_id: str):
        """Get or create a ChromaDB collection for a profile."""
        collection_name = f"profile_{profile_id.replace('-', '_')}"
        return self.client.get_or_create_collection(
            name=collection_name,
            metadata={"profile_id": profile_id}
        )
    
    def _get_tiered_collection(self, tier: str, tier_id: str):
        """
        Get or create a ChromaDB collection for a specific tier.
        
        Tiers:
        - 'user': User-level knowledge (My Universal Documents)
        - 'company': Company-level knowledge (Company-Wide KB)
        - 'profile': Profile-level knowledge (Strategic Profile KB) - existing behavior
        """
        # Sanitize the tier_id for collection name
        safe_id = tier_id.replace('-', '_')
        collection_name = f"{tier}_{safe_id}"
        
        return self.client.get_or_create_collection(
            name=collection_name,
            metadata={"tier": tier, "tier_id": tier_id}
        )
    
    async def process_document_tiered(
        self,
        file_path: str,
        tier: str,
        tier_id: str,
        filename: str,
        user_id: str
    ) -> Dict[str, Any]:
        """
        Process a document for a specific tier: extract text, chunk, embed, and store.
        
        Args:
            file_path: Path to the uploaded file
            tier: 'user', 'company', or 'profile'
            tier_id: user_id, company_id, or profile_id depending on tier
            filename: Original filename
            user_id: ID of the user uploading the document
        """
        document_id = str(uuid4())
        
        try:
            # Step 1: Extract text
            logger.info(f"Extracting text from {filename} for {tier} tier")
            text = await self._extract_text(file_path)
            
            if not text or len(text.strip()) < 10:
                raise ValueError("Document contains no extractable text")
            
            # Step 2: Chunk the text
            chunks = self._chunk_text(text)
            logger.info(f"Created {len(chunks)} chunks for {tier} tier document")
            
            # Step 3: Store chunks in ChromaDB
            collection = self._get_tiered_collection(tier, tier_id)
            
            chunk_ids = [f"{document_id}_chunk_{i}" for i in range(len(chunks))]
            chunk_metadatas = [
                {
                    "document_id": document_id,
                    "filename": filename,
                    "chunk_index": i,
                    "tier": tier,
                    "tier_id": tier_id,
                    "user_id": user_id
                }
                for i in range(len(chunks))
            ]
            
            collection.add(
                documents=chunks,
                ids=chunk_ids,
                metadatas=chunk_metadatas
            )
            
            # Step 4: Store document metadata in MongoDB
            doc_metadata = {
                "id": document_id,
                "tier": tier,
                "tier_id": tier_id,
                # Set the appropriate ID field based on tier
                f"{tier}_id": tier_id,
                "user_id": user_id,
                "filename": filename,
                "file_path": file_path,
                "file_size": Path(file_path).stat().st_size,
                "text_length": len(text),
                "chunk_count": len(chunks),
                "status": "processed",
                "created_at": datetime.now(timezone.utc).isoformat(),
                "processed_at": datetime.now(timezone.utc).isoformat()
            }
            
            if self.db is not None:
                await self.db.knowledge_documents.insert_one(doc_metadata)
            
            return {
                "success": True,
                "document_id": document_id,
                "filename": filename,
                "tier": tier,
                "text_length": len(text),
                "chunk_count": len(chunks),
                "status": "processed"
            }
            
        except Exception as e:
            logger.error(f"Error processing {tier} document {filename}: {str(e)}")
            
            # Store failed document metadata
            if self.db is not None:
                await self.db.knowledge_documents.insert_one({
                    "id": document_id,
                    "tier": tier,
                    "tier_id": tier_id,
                    f"{tier}_id": tier_id,
                    "user_id": user_id,
                    "filename": filename,
                    "file_path": file_path,
                    "status": "failed",
                    "error": str(e),
                    "created_at": datetime.now(timezone.utc).isoformat()
                })
            
            return {
                "success": False,
                "document_id": document_id,
                "filename": filename,
                "tier": tier,
                "error": str(e),
                "status": "failed"
            }
    
    async def query_tiered_knowledge_base(
        self,
        query: str,
        user_id: str,
        company_id: Optional[str] = None,
        profile_id: Optional[str] = None,
        n_results_per_tier: int = 3,
        profile_type: str = "personal"
    ) -> Dict[str, List[Dict[str, Any]]]:
        """
        Query knowledge bases across tiers based on profile type.
        
        FOUR-TIER Context-Aware Selection:
        - profile_type="company": Universal → Professional → User → Profile (4 sources)
        - profile_type="personal": Universal → User → Profile (3 sources, SKIPS Professional)
        
        This ensures:
        1. Universal Company Policies ALWAYS apply (Code of Conduct, Social Media Policy, etc.)
        2. Professional Brand & Compliance only applies to official company posts
        3. Personal content has creative freedom while respecting core company policies
        
        Args:
            query: The search query
            user_id: User's ID (for User-Level KB - Personal Compliance Rules)
            company_id: Company ID if user belongs to one (for Company-Level KBs)
            profile_id: Strategic Profile ID if selected (for Profile-Level KB)
            n_results_per_tier: Number of results to retrieve per tier
            profile_type: "personal" or "company" - determines which tiers to query
        """
        results = {
            "company_universal": [],    # Tier 1 - HIGHEST PRIORITY - Applies to ALL posts
            "company_professional": [], # Tier 2 - Professional Brand (only for company/brand profiles)
            "user": [],                 # Tier 3 - Personal compliance rules
            "profile": []               # Tier 4 - Strategic profile specific rules
        }
        
        # Query Universal Company Policies (Tier 1) - ALWAYS applies to ALL posts
        # This contains Code of Conduct, Acceptable Use Policy, etc.
        if company_id:
            try:
                collection = self._get_tiered_collection("company_universal", company_id)
                if collection.count() > 0:
                    universal_results = collection.query(
                        query_texts=[query],
                        n_results=min(n_results_per_tier, collection.count())
                    )
                    if universal_results and universal_results["documents"]:
                        for i, doc in enumerate(universal_results["documents"][0]):
                            results["company_universal"].append({
                                "content": doc,
                                "metadata": universal_results["metadatas"][0][i] if universal_results["metadatas"] else {},
                                "distance": universal_results["distances"][0][i] if universal_results["distances"] else None,
                                "tier": "company_universal"
                            })
                    logger.info(f"Retrieved {len(results['company_universal'])} universal company policy chunks (applies to ALL posts)")
            except Exception as e:
                logger.warning(f"Error querying universal company KB: {str(e)}")
        
        # Query Professional Brand & Compliance (Tier 2) - ONLY if profile_type is "company"
        # This contains brand guidelines, tone, product messaging - skipped for personal posts
        if company_id and profile_type == "company":
            try:
                collection = self._get_tiered_collection("company_professional", company_id)
                if collection.count() > 0:
                    professional_results = collection.query(
                        query_texts=[query],
                        n_results=min(n_results_per_tier, collection.count())
                    )
                    if professional_results and professional_results["documents"]:
                        for i, doc in enumerate(professional_results["documents"][0]):
                            results["company_professional"].append({
                                "content": doc,
                                "metadata": professional_results["metadatas"][0][i] if professional_results["metadatas"] else {},
                                "distance": professional_results["distances"][0][i] if professional_results["distances"] else None,
                                "tier": "company_professional"
                            })
                    logger.info(f"Retrieved {len(results['company_professional'])} professional brand chunks (profile_type: {profile_type})")
            except Exception as e:
                logger.warning(f"Error querying professional brand KB: {str(e)}")
        elif company_id and profile_type == "personal":
            logger.info(f"Skipping professional brand KB for personal profile (profile_type: {profile_type})")
        
        # Query User-Level KB (Tier 3) - Always query for both profile types
        # This contains personal compliance rules that apply to all content
        if user_id:
            try:
                collection = self._get_tiered_collection("user", user_id)
                if collection.count() > 0:
                    user_results = collection.query(
                        query_texts=[query],
                        n_results=min(n_results_per_tier, collection.count())
                    )
                    if user_results and user_results["documents"]:
                        for i, doc in enumerate(user_results["documents"][0]):
                            results["user"].append({
                                "content": doc,
                                "metadata": user_results["metadatas"][0][i] if user_results["metadatas"] else {},
                                "distance": user_results["distances"][0][i] if user_results["distances"] else None,
                                "tier": "user"
                            })
                    logger.info(f"Retrieved {len(results['user'])} user-level chunks")
            except Exception as e:
                logger.warning(f"Error querying user KB: {str(e)}")
        
        # Query Profile-Level KB (Tier 4) - existing behavior
        if profile_id:
            try:
                profile_results = await self.query_knowledge_base(profile_id, query, n_results_per_tier)
                for r in profile_results:
                    r["tier"] = "profile"
                results["profile"] = profile_results
                logger.info(f"Retrieved {len(results['profile'])} profile-level chunks")
            except Exception as e:
                logger.warning(f"Error querying profile KB: {str(e)}")
        
        return results
    
    async def get_tiered_context_for_ai(
        self,
        query: str,
        user_id: str,
        company_id: Optional[str] = None,
        profile_id: Optional[str] = None,
        profile_type: str = "personal"
    ) -> str:
        """
        Get combined context from tiers for AI content generation.
        
        FOUR-TIER CONTEXT-AWARE LOGIC:
        - profile_type="company": Universal → Professional → User → Profile (4 sources)
        - profile_type="personal": Universal → User → Profile (3 sources, SKIP Professional)
        
        This ensures:
        1. Universal Company Policies ALWAYS apply (Code of Conduct, etc.)
        2. Professional Brand & Compliance only applies to official company posts
        3. Personal content has creative freedom while respecting core policies
        
        Formats the context with clear tier labels and priority ordering.
        IMPORTANT: These are MANDATORY DIRECTIVES that the AI must follow.
        """
        tiered_results = await self.query_tiered_knowledge_base(
            query=query,
            user_id=user_id,
            company_id=company_id,
            profile_id=profile_id,
            profile_type=profile_type
        )
        
        context_parts = []
        
        # Determine preamble based on profile type
        if profile_type == "company":
            preamble = """⚠️ MANDATORY BRAND DIRECTIVES - YOU MUST FOLLOW THESE RULES ⚠️
The following guidelines are NON-NEGOTIABLE requirements from the user's brand documents.
These directives OVERRIDE your default behaviors and training.
Failure to follow these rules is a CRITICAL ERROR.
"""
        else:
            preamble = """⚠️ PERSONAL CONTENT GUIDELINES - FOLLOW THESE RULES ⚠️
The following are your personal content guidelines and preferences.
Note: Universal Company Policies still apply to protect the company.
"""
        
        # TIER 1: Universal Company Policies (HIGHEST PRIORITY) - ALWAYS applies to ALL posts
        # This contains Code of Conduct, Social Media Policy, Acceptable Use, etc.
        if tiered_results["company_universal"]:
            universal_context = "\n---\n".join([r["content"] for r in tiered_results["company_universal"]])
            context_parts.append(f"""🔴 UNIVERSAL COMPANY POLICIES (HIGHEST PRIORITY - APPLIES TO ALL POSTS):
These core company policies MUST be followed exactly. No exceptions allowed.
This includes Code of Conduct, Social Media Policy, and Acceptable Use guidelines.

{universal_context}""")
        
        # TIER 2: Professional Brand & Compliance - ONLY for company/brand profiles
        # This contains brand guidelines, tone, product messaging
        if tiered_results["company_professional"] and profile_type == "company":
            professional_context = "\n---\n".join([r["content"] for r in tiered_results["company_professional"]])
            context_parts.append(f"""🟣 PROFESSIONAL BRAND & COMPLIANCE (HIGH PRIORITY - COMPANY/BRAND POSTS ONLY):
These professional brand guidelines apply to official company communications.
Follow brand voice, tone, and messaging requirements.

{professional_context}""")
        
        # TIER 3: User-Level Context (Personal Compliance Rules) - Always included
        if tiered_results["user"]:
            user_context = "\n---\n".join([r["content"] for r in tiered_results["user"]])
            priority_label = "MEDIUM PRIORITY" if profile_type == "company" else "HIGH PRIORITY"
            context_parts.append(f"""🟠 PERSONAL COMPLIANCE RULES ({priority_label}):
Follow these personal preferences and constraints:

{user_context}""")
        
        # TIER 4: Profile-Level Context - Always included
        if tiered_results["profile"]:
            profile_context = "\n---\n".join([r["content"] for r in tiered_results["profile"]])
            context_parts.append(f"""🟡 PROFILE/CAMPAIGN STRATEGY (IMPORTANT):
Apply these strategic guidelines for this specific profile:

{profile_context}""")
        
        if not context_parts:
            return ""
        
        # Combine with preamble and closing reminder
        if profile_type == "company":
            closing = """
⚠️ REMINDER: The above directives are MANDATORY. Before generating content, verify that your output:
1. Does NOT violate any UNIVERSAL COMPANY POLICIES
2. Adheres to PROFESSIONAL BRAND guidelines for official communications
3. Follows ALL formatting requirements (currencies, units, terminology)
4. Adheres to brand voice and tone constraints
5. Avoids any prohibited topics, phrases, or formats mentioned above
"""
        else:
            closing = """
⚠️ REMINDER: Before generating content, verify that:
1. Your content respects UNIVERSAL COMPANY POLICIES (non-negotiable)
2. Your personal compliance rules are followed
3. The tone matches this specific profile's voice
4. Your authentic personal style is maintained while respecting core policies
"""
        
        logger.info(f"Generated AI context for profile_type={profile_type}, tiers used: universal={len(tiered_results['company_universal'])}, professional={len(tiered_results['company_professional'])}, user={len(tiered_results['user'])}, profile={len(tiered_results['profile'])}")
        
        return preamble + "\n\n".join(context_parts) + closing
    
    async def delete_document_tiered(
        self,
        document_id: str,
        tier: str,
        tier_id: str
    ) -> bool:
        """Delete a document from a tiered knowledge base."""
        try:
            # Delete from ChromaDB
            collection = self._get_tiered_collection(tier, tier_id)
            
            results = collection.get(
                where={"document_id": document_id}
            )
            
            if results and results["ids"]:
                collection.delete(ids=results["ids"])
                logger.info(f"Deleted {len(results['ids'])} chunks for {tier} document {document_id}")
            
            # Delete from MongoDB
            if self.db is not None:
                doc = await self.db.knowledge_documents.find_one({"id": document_id})
                if doc and doc.get("file_path"):
                    file_path = Path(doc["file_path"])
                    if file_path.exists():
                        file_path.unlink()
                
                await self.db.knowledge_documents.delete_one({"id": document_id})
            
            return True
            
        except Exception as e:
            logger.error(f"Error deleting {tier} document {document_id}: {str(e)}")
            return False
    
    async def get_tiered_stats(self, tier: str, tier_id: str) -> Dict[str, Any]:
        """Get statistics for a specific tier's knowledge base."""
        try:
            collection = self._get_tiered_collection(tier, tier_id)
            chunk_count = collection.count()
            
            doc_count = 0
            if self.db is not None:
                doc_count = await self.db.knowledge_documents.count_documents({
                    "tier": tier,
                    "tier_id": tier_id,
                    "status": "processed"
                })
            
            return {
                "document_count": doc_count,
                "chunk_count": chunk_count,
                "has_knowledge": chunk_count > 0
            }
        except Exception as e:
            logger.error(f"Error getting {tier} stats: {str(e)}")
            return {"document_count": 0, "chunk_count": 0, "has_knowledge": False}

    async def process_document(
        self,
        file_path: str,
        profile_id: str,
        filename: str,
        user_id: str
    ) -> Dict[str, Any]:
        """
        Process a document: extract text, chunk, embed, and store in vector DB.
        This is the core RAG preprocessing step.
        """
        document_id = str(uuid4())
        
        try:
            # Step 1: Extract text from document
            logger.info(f"Extracting text from {filename}")
            text = await self._extract_text(file_path)
            
            if not text or len(text.strip()) < 10:
                raise ValueError("Document contains no extractable text")
            
            # Step 2: Chunk the text
            logger.info(f"Chunking text ({len(text)} chars)")
            chunks = self._chunk_text(text)
            logger.info(f"Created {len(chunks)} chunks")
            
            # Step 3: Store chunks in ChromaDB (it handles embeddings automatically)
            collection = self._get_collection(profile_id)
            
            # Prepare data for ChromaDB
            ids = [f"{document_id}_{i}" for i in range(len(chunks))]
            metadatas = [
                {
                    "document_id": document_id,
                    "profile_id": profile_id,
                    "filename": filename,
                    "chunk_index": i,
                    "total_chunks": len(chunks)
                }
                for i in range(len(chunks))
            ]
            
            # Add to collection (ChromaDB embeds automatically using default model)
            collection.add(
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )
            
            logger.info(f"Stored {len(chunks)} chunks in ChromaDB for profile {profile_id}")
            
            # Step 4: Store document metadata in MongoDB
            doc_metadata = {
                "id": document_id,
                "profile_id": profile_id,
                "user_id": user_id,
                "filename": filename,
                "file_path": file_path,
                "file_size": os.path.getsize(file_path),
                "text_length": len(text),
                "chunk_count": len(chunks),
                "status": "processed",
                "created_at": datetime.now(timezone.utc).isoformat(),
                "processed_at": datetime.now(timezone.utc).isoformat()
            }
            
            if self.db is not None:
                await self.db.knowledge_documents.insert_one(doc_metadata)
            
            return {
                "success": True,
                "document_id": document_id,
                "filename": filename,
                "text_length": len(text),
                "chunk_count": len(chunks),
                "status": "processed"
            }
            
        except Exception as e:
            logger.error(f"Error processing document {filename}: {str(e)}")
            
            # Store failed document metadata
            if self.db is not None:
                await self.db.knowledge_documents.insert_one({
                    "id": document_id,
                    "profile_id": profile_id,
                    "user_id": user_id,
                    "filename": filename,
                    "file_path": file_path,
                    "status": "failed",
                    "error": str(e),
                    "created_at": datetime.now(timezone.utc).isoformat()
                })
            
            return {
                "success": False,
                "document_id": document_id,
                "filename": filename,
                "error": str(e),
                "status": "failed"
            }
    
    async def _extract_text(self, file_path: str) -> str:
        """Extract text from various document formats."""
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        
        try:
            if suffix == ".pdf":
                return self._extract_pdf(file_path)
            elif suffix in [".docx", ".doc"]:
                return self._extract_docx(file_path)
            elif suffix in [".xlsx", ".xls"]:
                return self._extract_xlsx(file_path)
            elif suffix in [".pptx", ".ppt"]:
                return self._extract_pptx(file_path)
            elif suffix in [".txt", ".md", ".csv"]:
                return self._extract_text_file(file_path)
            else:
                raise ValueError(f"Unsupported file format: {suffix}")
        except Exception as e:
            logger.error(f"Text extraction error for {file_path}: {str(e)}")
            raise
    
    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files."""
        text_parts = []
        reader = PdfReader(str(file_path))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        return "\n\n".join(text_parts)
    
    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from Word documents."""
        doc = DocxDocument(str(file_path))
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        return "\n\n".join(text_parts)
    
    def _extract_xlsx(self, file_path: Path) -> str:
        """Extract text from Excel files."""
        wb = load_workbook(str(file_path), data_only=True)
        text_parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text_parts.append(f"Sheet: {sheet}")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts)
    
    def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint files."""
        prs = Presentation(str(file_path))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {slide_num}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text_parts.append("\n".join(slide_text))
        return "\n\n".join(text_parts)
    
    def _extract_text_file(self, file_path: Path) -> str:
        """Extract text from plain text files."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks for better retrieval."""
        chunks = []
        
        # Clean the text
        text = " ".join(text.split())
        
        if len(text) <= CHUNK_SIZE:
            return [text]
        
        start = 0
        while start < len(text):
            end = start + CHUNK_SIZE
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence ending punctuation
                for punct in [". ", "! ", "? ", "\n"]:
                    last_punct = text.rfind(punct, start, end)
                    if last_punct > start + CHUNK_SIZE // 2:
                        end = last_punct + 1
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - CHUNK_OVERLAP
        
        return chunks
    
    async def query_knowledge_base(
        self,
        profile_id: str,
        query: str,
        n_results: int = 5
    ) -> List[Dict[str, Any]]:
        """
        Query the knowledge base for relevant context.
        This is the RAG retrieval step.
        """
        try:
            collection = self._get_collection(profile_id)
            
            # Check if collection has documents
            if collection.count() == 0:
                logger.info(f"No documents in knowledge base for profile {profile_id}")
                return []
            
            # Query ChromaDB (it handles embedding the query automatically)
            results = collection.query(
                query_texts=[query],
                n_results=min(n_results, collection.count())
            )
            
            # Format results
            formatted_results = []
            if results and results["documents"]:
                for i, doc in enumerate(results["documents"][0]):
                    formatted_results.append({
                        "content": doc,
                        "metadata": results["metadatas"][0][i] if results["metadatas"] else {},
                        "distance": results["distances"][0][i] if results["distances"] else None
                    })
            
            logger.info(f"Retrieved {len(formatted_results)} relevant chunks for profile {profile_id}")
            return formatted_results
            
        except Exception as e:
            logger.error(f"Error querying knowledge base: {str(e)}")
            return []
    
    async def delete_document(self, document_id: str, profile_id: str) -> bool:
        """Delete a document and its chunks from the knowledge base."""
        try:
            # Delete from ChromaDB
            collection = self._get_collection(profile_id)
            
            # Get all chunk IDs for this document
            results = collection.get(
                where={"document_id": document_id}
            )
            
            if results and results["ids"]:
                collection.delete(ids=results["ids"])
                logger.info(f"Deleted {len(results['ids'])} chunks for document {document_id}")
            
            # Delete from MongoDB
            if self.db is not None:
                doc = await self.db.knowledge_documents.find_one({"id": document_id})
                if doc and doc.get("file_path"):
                    # Delete the file
                    file_path = Path(doc["file_path"])
                    if file_path.exists():
                        file_path.unlink()
                
                await self.db.knowledge_documents.delete_one({"id": document_id})
            
            return True
            
        except Exception as e:
            logger.error(f"Error deleting document {document_id}: {str(e)}")
            return False
    
    async def get_profile_documents(self, profile_id: str) -> List[Dict]:
        """Get all documents for a profile."""
        if self.db is None:
            return []
        
        docs = await self.db.knowledge_documents.find(
            {"profile_id": profile_id},
            {"_id": 0}
        ).to_list(100)
        
        return docs
    
    async def get_profile_stats(self, profile_id: str) -> Dict[str, Any]:
        """Get statistics about a profile's knowledge base."""
        try:
            collection = self._get_collection(profile_id)
            chunk_count = collection.count()
            
            doc_count = 0
            if self.db is not None:
                doc_count = await self.db.knowledge_documents.count_documents(
                    {"profile_id": profile_id, "status": "processed"}
                )
            
            return {
                "document_count": doc_count,
                "chunk_count": chunk_count,
                "has_knowledge": chunk_count > 0
            }
        except Exception as e:
            logger.error(f"Error getting profile stats: {str(e)}")
            return {"document_count": 0, "chunk_count": 0, "has_knowledge": False}
    
    async def get_knowledge_summary(self, profile_id: str, max_chunks: int = 10) -> str:
        """
        Get a representative summary of the knowledge base for SEO keyword generation.
        Retrieves diverse chunks from the knowledge base to capture key themes.
        """
        try:
            collection = self._get_collection(profile_id)
            
            # Check if collection has documents
            total_chunks = collection.count()
            if total_chunks == 0:
                logger.info(f"No knowledge base content for profile {profile_id}")
                return ""
            
            # Get a sample of chunks for diversity
            # Use broad queries to capture different themes
            broad_queries = [
                "main products services offerings",
                "company business description overview",
                "target audience customers market",
                "unique features benefits advantages",
                "industry sector expertise specialization"
            ]
            
            all_chunks = []
            seen_contents = set()
            
            for query in broad_queries:
                if len(all_chunks) >= max_chunks:
                    break
                    
                try:
                    results = collection.query(
                        query_texts=[query],
                        n_results=min(3, total_chunks)
                    )
                    
                    if results and results["documents"]:
                        for doc in results["documents"][0]:
                            # Deduplicate chunks
                            content_hash = hash(doc[:100])
                            if content_hash not in seen_contents and len(all_chunks) < max_chunks:
                                seen_contents.add(content_hash)
                                all_chunks.append(doc)
                except Exception as e:
                    logger.warning(f"Query failed for '{query}': {str(e)}")
                    continue
            
            # If we still don't have enough, try to get some random samples
            if len(all_chunks) < max_chunks // 2:
                try:
                    # Get all chunks and sample
                    all_results = collection.get(limit=min(50, total_chunks))
                    if all_results and all_results["documents"]:
                        for doc in all_results["documents"]:
                            content_hash = hash(doc[:100])
                            if content_hash not in seen_contents and len(all_chunks) < max_chunks:
                                seen_contents.add(content_hash)
                                all_chunks.append(doc)
                except Exception as e:
                    logger.warning(f"Failed to get additional chunks: {str(e)}")
            
            # Combine chunks into a summary (limit total length)
            combined_text = "\n\n---\n\n".join(all_chunks)
            
            # Truncate if too long (keep under 4000 chars for API efficiency)
            if len(combined_text) > 4000:
                combined_text = combined_text[:4000] + "..."
            
            logger.info(f"Generated knowledge summary ({len(combined_text)} chars) from {len(all_chunks)} chunks for profile {profile_id}")
            return combined_text
            
        except Exception as e:
            logger.error(f"Error generating knowledge summary: {str(e)}")
            return ""
    
    async def get_document_chunks(self, document_id: str) -> List[Dict[str, Any]]:
        """
        Get all text chunks extracted from a specific document.
        Returns chunks with their content for user transparency.
        """
        try:
            # Get document info
            doc = await self.db.knowledge_documents.find_one(
                {"id": document_id},
                {"_id": 0}
            )
            
            if not doc:
                logger.warning(f"Document {document_id} not found")
                return []
            
            profile_id = doc.get("profile_id")
            
            # Get the collection for this profile
            collection = self._get_collection(profile_id)
            
            # Query all chunks for this document
            # ChromaDB stores document chunks with metadata containing the document_id
            all_data = collection.get(
                where={"document_id": document_id}
            )
            
            if not all_data or not all_data.get("documents"):
                # Try alternative: get all and filter
                all_data = collection.get()
                chunks = []
                
                if all_data and all_data.get("ids"):
                    for i, chunk_id in enumerate(all_data["ids"]):
                        # Check if this chunk belongs to the document
                        if chunk_id.startswith(document_id):
                            content = all_data["documents"][i] if all_data.get("documents") else ""
                            chunks.append({
                                "id": chunk_id,
                                "content": content
                            })
                
                logger.info(f"Retrieved {len(chunks)} chunks for document {document_id}")
                return chunks
            
            # Build chunks list
            chunks = []
            for i, doc_content in enumerate(all_data.get("documents", [])):
                chunk_id = all_data["ids"][i] if all_data.get("ids") else f"chunk_{i}"
                chunks.append({
                    "id": chunk_id,
                    "content": doc_content
                })
            
            logger.info(f"Retrieved {len(chunks)} chunks for document {document_id}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error getting document chunks: {str(e)}")
            return []
    
    async def get_context_for_generation(self, profile_id: str, limit: int = 5) -> str:
        """
        Get context from knowledge base for AI content generation.
        Returns a string of relevant chunks.
        """
        try:
            collection = self._get_collection(profile_id)
            
            # Get some representative chunks
            results = collection.get(limit=limit)
            
            if not results or not results.get("documents"):
                return ""
            
            # Combine chunks
            chunks = results["documents"]
            combined = "\n\n---\n\n".join(chunks)
            
            # Truncate if too long
            if len(combined) > 3000:
                combined = combined[:3000] + "..."
            
            return combined
            
        except Exception as e:
            logger.error(f"Error getting context for generation: {str(e)}")
            return ""


# Global instance
_knowledge_service: Optional[KnowledgeBaseService] = None


def init_knowledge_service(db) -> KnowledgeBaseService:
    """Initialize the knowledge base service."""
    global _knowledge_service
    _knowledge_service = KnowledgeBaseService(db)
    logger.info("Knowledge Base Service initialized")
    return _knowledge_service


def get_knowledge_service() -> KnowledgeBaseService:
    """Get the knowledge base service instance."""
    if _knowledge_service is None:
        raise RuntimeError("Knowledge service not initialized. Call init_knowledge_service first.")
    return _knowledge_service
